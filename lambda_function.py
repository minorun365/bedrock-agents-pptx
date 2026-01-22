import json
import os
import boto3
import urllib.request
from io import BytesIO
from datetime import datetime

# python-pptxはLambdaレイヤーとして追加する必要あり
from pptx import Presentation
from pptx.util import Inches, Pt


def lambda_handler(event, context):
    """Bedrock Agentsからのリクエストを処理"""
    print(f"Received event: {json.dumps(event)}")

    action_group = event.get("actionGroup", "")
    function_name = event.get("function", "")
    parameters = {p["name"]: p["value"] for p in event.get("parameters", [])}

    # 機能ごとに処理を分岐
    if function_name == "search-web":
        result = search_web(parameters.get("query", ""))
    elif function_name == "create-pptx":
        result = create_pptx(
            title=parameters.get("title", "無題"),
            content=parameters.get("content", "")
        )
    elif function_name == "send-email":
        result = send_email(url=parameters.get("url", ""))
    else:
        result = {"error": f"Unknown function: {function_name}"}

    # Bedrock Agents用のレスポンス形式
    return {
        "messageVersion": "1.0",
        "response": {
            "actionGroup": action_group,
            "function": function_name,
            "functionResponse": {
                "responseBody": {
                    "TEXT": {
                        "body": json.dumps(result, ensure_ascii=False)
                    }
                }
            }
        }
    }


def search_web(query: str) -> dict:
    """Tavily APIを使用してWeb検索"""
    api_key = os.environ.get("TAVILY_API_KEY")
    if not api_key:
        return {"error": "TAVILY_API_KEY not set"}

    url = "https://api.tavily.com/search"
    data = json.dumps({
        "api_key": api_key,
        "query": query,
        "max_results": 5
    }).encode("utf-8")

    req = urllib.request.Request(
        url,
        data=data,
        headers={"Content-Type": "application/json"}
    )

    try:
        with urllib.request.urlopen(req, timeout=30) as res:
            response = json.loads(res.read().decode("utf-8"))
            # 検索結果を整形
            results = []
            for r in response.get("results", []):
                results.append({
                    "title": r.get("title", ""),
                    "url": r.get("url", ""),
                    "content": r.get("content", "")
                })
            return {"query": query, "results": results}
    except Exception as e:
        return {"error": str(e)}


def create_pptx(title: str, content: str) -> dict:
    """PowerPointファイルを作成してS3にアップロード"""
    bucket = os.environ.get("S3_BUCKET")
    if not bucket:
        return {"error": "S3_BUCKET not set"}

    # プレゼンテーション作成
    prs = Presentation()

    # タイトルスライド
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"

    # コンテンツを\n\nで分割して複数スライド作成
    content = content.strip()
    slides_content = content.split('\n\n') if content else []

    for slide_content in slides_content:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)

        lines = slide_content.split('\n')
        # 1行目をタイトルに
        slide.shapes.title.text = lines[0].lstrip('- ').lstrip('# ')
        # 2行目以降を本文に
        if len(lines) > 1:
            body = slide.placeholders[1]
            body.text = '\n'.join([line.lstrip('- ') for line in lines[1:]])

    # BytesIOに保存
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)

    # S3にアップロード
    s3 = boto3.client("s3")
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    file_key = f"slide_{timestamp}.pptx"

    try:
        s3.upload_fileobj(pptx_buffer, bucket, file_key)
        # 署名付きURLを生成（1時間有効）
        presigned_url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": bucket, "Key": file_key},
            ExpiresIn=3600
        )
        return {
            "message": "PowerPoint created successfully",
            "s3_key": file_key,
            "download_url": presigned_url
        }
    except Exception as e:
        return {"error": str(e)}


def send_email(url: str) -> dict:
    """Amazon SNSでメール送信"""
    topic_arn = os.environ.get("SNS_TOPIC_ARN")
    if not topic_arn:
        return {"error": "SNS_TOPIC_ARN not set"}

    sns = boto3.client("sns")

    try:
        response = sns.publish(
            TopicArn=topic_arn,
            Subject="Bedrock Agentsがパワポを作成しました",
            Message=f"以下のURLからダウンロードできます：\n{url}"
        )
        return {
            "message": "Email sent successfully",
            "message_id": response.get("MessageId")
        }
    except Exception as e:
        return {"error": str(e)}
